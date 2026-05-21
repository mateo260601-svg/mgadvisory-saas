from pathlib import Path

from app.services.template_service import load_pptx_template_manifest


BRAND_GREEN = "4F8758"
BRAND_DARK = "252B2A"
GRAPHITE = "3F4A47"
LIGHT_BG = "EEF1EF"
WHITE = "FFFFFF"
MUTED = "68726F"


def build_lender_presentation(project: dict, financials: dict, output_path: Path, blueprint: dict | None = None, assumptions: dict | None = None) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:
        raise RuntimeError(f"PPTX generation dependency unavailable: {exc}") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ctx = {
        "Presentation": Presentation,
        "RGBColor": RGBColor,
        "CategoryChartData": CategoryChartData,
        "XL_CHART_TYPE": XL_CHART_TYPE,
        "XL_LABEL_POSITION": XL_LABEL_POSITION,
        "XL_LEGEND_POSITION": XL_LEGEND_POSITION,
        "MSO_SHAPE": MSO_SHAPE,
        "PP_ALIGN": PP_ALIGN,
        "Inches": Inches,
        "Pt": Pt,
        "project": project,
        "financials": financials,
        "template_manifest": load_pptx_template_manifest(),
        "blueprint": blueprint or {},
        "assumptions": assumptions or {},
    }

    _cover(prs, ctx)
    _executive_snapshot(prs, ctx)
    _historical_financials(prs, ctx)
    _bp_output_charts(prs, ctx)
    _revenue_split(prs, ctx)
    _business_plan_bridge(prs, ctx)
    _debt_and_covenants(prs, ctx)
    _capital_structure(prs, ctx)
    _restructuring_options(prs, ctx)
    _key_diligence(prs, ctx)
    _claude_slide_blueprint(prs, ctx)
    _template_standard(prs, ctx)
    _appendix_model_map(prs, ctx)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def _cover(prs, ctx):
    slide = _blank_slide(prs, ctx, dark=True)
    company = ctx["project"].get("company_name", "Target Company")
    _text(slide, ctx, 0.7, 0.65, 8.1, 0.35, "MG STRATEGIC FINANCE AI", size=10, color=LIGHT_BG, bold=True)
    _text(slide, ctx, 0.7, 1.52, 8.2, 1.1, company, size=34, color=WHITE, bold=True)
    _text(slide, ctx, 0.72, 2.75, 6.4, 0.7, "Institutional lender presentation | Business plan, debt capacity and restructuring options", size=16, color="DCE4DF")
    _metric_band(
        slide,
        ctx,
        0.72,
        4.55,
        [
            ("Project", ctx["project"].get("project_type", "Investment case")),
            ("Currency", ctx["project"].get("currency", "EUR")),
            ("Model", "60-month BP"),
            ("Debt", "10-tranche engine"),
        ],
        dark=True,
    )
    _footer(slide, ctx, "Confidential | Draft for discussion")


def _executive_snapshot(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Executive snapshot", "The model translates uploaded historicals into a lender-ready planning and credit analysis pack.")
    metrics = _summary_metrics(ctx["financials"])
    _metric_band(slide, ctx, 0.6, 1.55, metrics)
    _section_label(slide, ctx, 0.6, 3.05, "Credit read-across")
    bullets = [
        "Revenue, EBITDA, cash flow and leverage are linked through a formula-driven monthly model.",
        "Debt schedule captures start dates, maturities, moratoriums, PIK, bullet repayment and cash sweep mechanics.",
        "Covenant headroom, liquidity and refinancing risk are packaged into dedicated output tabs.",
        "AI extraction creates a traceable first pass from uploaded PDF/XLSX/CSV documents, with confidence flags.",
    ]
    _bullets(slide, ctx, 0.75, 3.42, 6.6, 2.4, bullets)
    _callout(slide, ctx, 8.35, 3.18, 4.15, 2.4, "Analyst note", "Use this pack as a first institutional draft. Final external release should include management validation of extracted historicals and covenant definitions.")
    _footer(slide, ctx, "Executive summary")


def _historical_financials(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Historical financials", "Extracted source data is standardised into P&L, balance sheet, cash flow, debt and working-capital views.")
    table_data = _historical_table(ctx["financials"])
    _table(slide, ctx, 0.65, 1.55, 7.25, 3.55, table_data)
    extraction = ctx["financials"].get("extraction", {})
    _callout(
        slide,
        ctx,
        8.25,
        1.55,
        4.3,
        1.6,
        "Extraction status",
        f"Mode: {extraction.get('mode', 'local')}\nConfidence: {extraction.get('confidence', 'n/a')}",
    )
    _callout(
        slide,
        ctx,
        8.25,
        3.45,
        4.3,
        1.9,
        "Required clean-up",
        "Map statutory line items to lender definitions, validate exceptional items, confirm units and reconcile EBITDA to management reporting.",
    )
    _footer(slide, ctx, "Historical extraction")


def _business_plan_bridge(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Business plan bridge", "Forecast outputs are built from explicit drivers, not hardcoded presentation numbers.")
    _section_label(slide, ctx, 0.65, 1.38, "Driver architecture")
    steps = [
        ("01", "Historical inputs", "Uploaded documents and AI extraction feed normalised historicals."),
        ("02", "Commercial drivers", "Product/service volumes, pricing and growth drive revenue."),
        ("03", "Cost base", "COGS, payroll, opex, working capital and capex flow into statements."),
        ("04", "Debt mechanics", "Tranche-level debt schedule calculates interest, PIK, amortisation and sweeps."),
        ("05", "Outputs", "Covenants, debt capacity, sensitivities and IC summary are linked to the model."),
    ]
    for i, (num, title, body) in enumerate(steps):
        _step(slide, ctx, 0.75 + i * 2.45, 2.05, num, title, body)
    _footer(slide, ctx, "Business plan model")


def _bp_output_charts(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "BP output dashboard", "The deck converts the model into banker-readable revenue, EBITDA and leverage proof points.")
    years, revenue, ebitda, debt = _bp_projection_series(ctx)
    _chart(slide, ctx, 0.65, 1.65, 5.85, 3.25, "Revenue / EBITDA", years, [("Revenue", revenue), ("EBITDA", ebitda)], chart_type="column")
    leverage = [(d / e if e else 0) for d, e in zip(debt, ebitda)]
    _chart(slide, ctx, 7.0, 1.65, 5.65, 3.25, "Net leverage", years, [("Net debt / EBITDA", leverage)], chart_type="line")
    _callout(slide, ctx, 0.65, 5.25, 5.85, 1.05, "Source", "Charts are generated from the BP Builder assumptions and refreshed when the deck is regenerated.")
    _callout(slide, ctx, 7.0, 5.25, 5.65, 1.05, "Review point", "For final banking materials, reconcile these charts to the downloadable Excel BP and management-approved cases.")
    _footer(slide, ctx, "BP outputs")


def _revenue_split(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Product and service revenue split", "Custom product/service names flow from the SaaS BP Builder into the deck narrative.")
    rows = [["Product / service", "Type", "Start volume", "Start price", "Implied monthly revenue"]]
    values = []
    labels = []
    for item in (ctx.get("assumptions", {}).get("revenue_streams") or [])[:5]:
        revenue = float(item.get("volume") or 0) * float(item.get("price") or 0)
        labels.append(item.get("name") or "Revenue stream")
        values.append(revenue)
        rows.append([item.get("name", ""), item.get("type", ""), f"{float(item.get('volume') or 0):,.0f}", _money(float(item.get("price") or 0)), _money(revenue)])
    if len(rows) == 1:
        rows.append(["No BP Builder revenue streams", "-", "-", "-", "-"])
        labels, values = ["No data"], [1]
    _table(slide, ctx, 0.65, 1.55, 7.0, 3.35, rows)
    _chart(slide, ctx, 8.0, 1.65, 4.4, 3.15, "Opening revenue mix", labels, [("Revenue", values)], chart_type="pie")
    _footer(slide, ctx, "Revenue split")


def _debt_and_covenants(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Debt and covenant analysis", "The lender view focuses on serviceability, covenant headroom, liquidity and refinancing risk.")
    _table(
        slide,
        ctx,
        0.65,
        1.55,
        7.25,
        3.85,
        [
            ["Analytical module", "Purpose", "Model output"],
            ["Debt schedule", "Tranche-level debt roll-forward", "Opening debt, draws, interest, PIK, amortisation, closing debt"],
            ["Covenants", "Test lender thresholds", "Net debt / EBITDA, ICR, liquidity, pass/fail"],
            ["Debt capacity", "Size supportable debt", "Implied debt capacity and headroom / shortfall"],
            ["Sensitivity matrix", "Downside stress", "Leverage and liquidity under EBITDA / debt / FCF shocks"],
        ],
    )
    _callout(slide, ctx, 8.25, 1.55, 4.25, 1.55, "Bank case lens", "Focus lender narrative on downside cash generation, cure capacity and the first covenant tripwire.")
    _callout(slide, ctx, 8.25, 3.45, 4.25, 1.95, "Model rule", "Every debt output is formula-linked to assumptions. No presentation number should be typed over.")
    _footer(slide, ctx, "Debt analytics")


def _capital_structure(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Capital structure and interest mechanics", "Debt instruments are selected manually by dossier, including cash pay / PIK and payment frequency.")
    rows = [["Instrument", "Type", "Opening", "Commitment", "Interest", "Frequency", "Amortisation"]]
    labels = []
    balances = []
    for item in (ctx.get("assumptions", {}).get("debt_tranches") or [])[:5]:
        opening = float(item.get("opening_balance") or 0)
        commitment = float(item.get("commitment") or 0)
        labels.append(item.get("name") or item.get("debt_type") or "Debt")
        balances.append(opening)
        rows.append([
            item.get("name", ""),
            item.get("debt_type", ""),
            _money(opening),
            _money(commitment),
            item.get("interest_type") or ("PIK" if item.get("pik") else "Cash"),
            item.get("cash_pay_frequency", "Monthly"),
            item.get("amortization", ""),
        ])
    if len(rows) == 1:
        rows.append(["No debt configured", "-", "-", "-", "-", "-", "-"])
        labels, balances = ["No debt"], [0]
    _table(slide, ctx, 0.55, 1.55, 8.0, 3.6, rows)
    _chart(slide, ctx, 8.9, 1.65, 3.75, 3.2, "Opening debt stack", labels, [("Opening debt", balances)], chart_type="bar")
    _callout(slide, ctx, 0.65, 5.45, 11.8, 0.85, "Model implication", "Cash interest affects liquidity and free cash flow, while PIK compounds closing debt and leverage until maturity or refinancing.")
    _footer(slide, ctx, "Capital structure")


def _restructuring_options(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Restructuring options paper", "Options are sequenced by liquidity relief, execution risk, lender acceptability and value preservation.")
    options = [
        ["Option", "Liquidity", "Complexity", "Use case"],
        ["Self-help", "Low", "Low", "Cost reduction, working-capital release"],
        ["Amend & extend", "Medium", "Medium", "Temporary covenant pressure"],
        ["PIK / payment holiday", "High", "Medium", "Short-term liquidity protection"],
        ["Super senior new money", "High", "High", "Liquidity runway insufficient"],
        ["Debt-for-equity", "Medium", "Very High", "Unsustainable capital structure"],
    ]
    _table(slide, ctx, 0.65, 1.55, 7.5, 3.95, options)
    _callout(slide, ctx, 8.45, 1.55, 3.95, 3.95, "Decision logic", "Start with the least value-destructive option that restores liquidity and covenant stability. Escalate only if the forecast shows unsustainable serviceability or insufficient stakeholder support.")
    _footer(slide, ctx, "Restructuring options")


def _key_diligence(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Priority diligence questions", "The next workstream should convert extraction output into validated lender-grade evidence.")
    boxes = [
        ("Quality of earnings", "Which EBITDA items are recurring, exceptional or owner-adjusted?"),
        ("Cash conversion", "Are DSO/DPO/inventory assumptions grounded in ageing and supplier behaviour?"),
        ("Debt documents", "Do covenant definitions match the modelled numerator and denominator?"),
        ("Liquidity runway", "What minimum cash is truly required to operate without supplier/customer disruption?"),
        ("Scenario governance", "Who owns Base, Bank Case, Downside and Restructuring Case assumptions?"),
        ("Data traceability", "Can every historical number be tied back to a source file and line item?"),
    ]
    for idx, (title, body) in enumerate(boxes):
        x = 0.65 + (idx % 3) * 4.18
        y = 1.55 + (idx // 3) * 2.05
        _callout(slide, ctx, x, y, 3.75, 1.45, title, body)
    _footer(slide, ctx, "Diligence workplan")


def _appendix_model_map(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Appendix: model architecture", "The workbook is organised as an auditable institutional finance model.")
    modules = [
        ["Area", "Key tabs"],
        ["Inputs", "Admin, Control Panel, Historical Inputs, Group Assumptions, Data Room"],
        ["Operations", "Revenue Drivers, Product Build, Headcount, Opex, Working Capital, Capex D&A"],
        ["Debt", "Debt Config, Debt Schedule, Covenants, Debt Capacity"],
        ["Outputs", "Financial Statements, Outputs, Packaged Output, IC Summary, Sensitivity Matrix"],
        ["Controls", "Lookup, Mapping, Checks, Lists & Dates"],
    ]
    _table(slide, ctx, 0.75, 1.55, 11.8, 3.95, modules)
    _footer(slide, ctx, "Appendix")


def _template_standard(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Template standard", "Generated decks are benchmarked against the institutional reference books loaded in the SaaS.")
    manifest = ctx.get("template_manifest", {})
    templates = manifest.get("templates", [])
    rows = [["Reference book", "Role", "Slides", "How it is used"]]
    for item in templates[:4]:
        filename = item.get("file", "")
        if "book_schemas" in filename:
            role = "Pattern library"
            use = "Charts, timelines, process pages, tables, dividers and analytics layouts"
        elif "im_draft" in filename:
            role = "IM narrative"
            use = "Executive flow, investment highlights, value creation and next steps"
        else:
            role = "Reference"
            use = "Design grammar and slide density"
        rows.append([filename.replace("_reference.pptx", ""), role, item.get("slide_count", ""), use])
    if len(rows) == 1:
        rows.append(["No template manifest loaded", "Fallback", "-", "Uses built-in MG institutional layout system"])
    _table(slide, ctx, 0.65, 1.55, 11.9, 2.55, rows)
    _callout(
        slide,
        ctx,
        0.65,
        4.55,
        5.8,
        1.35,
        "Production principle",
        "The generator inherits rhythm, density and proof-object discipline from the books, while replacing client-specific wording with the current project narrative.",
    )
    _callout(
        slide,
        ctx,
        6.8,
        4.55,
        5.75,
        1.35,
        "Quality target",
        "Every generated deck should read as a first institutional draft for bankers, lenders, restructuring advisers or transaction services teams.",
    )
    _footer(slide, ctx, "Template standard")


def _claude_slide_blueprint(prs, ctx):
    slide = _blank_slide(prs, ctx)
    _headline(slide, ctx, "Claude slide blueprint", "Claude selects the most relevant reference pattern for each planned slide before generation.")
    plan = (ctx.get("blueprint") or {}).get("blueprint", {})
    rows = [["#", "Generated slide", "Reference pattern", "Proof object"]]
    for item in plan.get("slides", [])[:7]:
        rows.append(
            [
                item.get("slide_number", ""),
                item.get("slide_title", ""),
                f"{item.get('recommended_template_file', '').replace('_reference.pptx', '')} / slide {item.get('recommended_template_slide', '')}",
                item.get("proof_object", ""),
            ]
        )
    if len(rows) == 1:
        rows.append(["-", "No Claude blueprint available", "Fallback MG layout", "mixed"])
    _table(slide, ctx, 0.55, 1.55, 12.15, 3.7, rows)
    thesis = plan.get("narrative_thesis") or "The deck uses fallback institutional layouts if Claude is not configured."
    _callout(slide, ctx, 0.65, 5.55, 11.8, 0.9, "Narrative thesis", thesis)
    _footer(slide, ctx, "Claude template selection")


def _blank_slide(prs, ctx, dark: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    RGBColor = ctx["RGBColor"]
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(BRAND_DARK if dark else LIGHT_BG)
    if not dark:
        shape = slide.shapes.add_shape(ctx["MSO_SHAPE"].RECTANGLE, ctx["Inches"](0), ctx["Inches"](0), ctx["Inches"](13.333), ctx["Inches"](0.12))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(BRAND_GREEN)
        shape.line.fill.background()
    return slide


def _headline(slide, ctx, title: str, subtitle: str):
    _text(slide, ctx, 0.6, 0.38, 8.6, 0.42, "MG STRATEGIC FINANCE AI", size=8, color=BRAND_GREEN, bold=True)
    _text(slide, ctx, 0.6, 0.78, 7.6, 0.5, title, size=25, color=BRAND_DARK, bold=True)
    _text(slide, ctx, 0.62, 1.22, 9.2, 0.35, subtitle, size=10, color=MUTED)


def _section_label(slide, ctx, x, y, label):
    _text(slide, ctx, x, y, 3.2, 0.25, label.upper(), size=8, color=BRAND_GREEN, bold=True)


def _metric_band(slide, ctx, x, y, metrics, dark: bool = False):
    for idx, (label, value) in enumerate(metrics):
        _callout(slide, ctx, x + idx * 2.95, y, 2.62, 0.95, label, str(value), dark=dark)


def _callout(slide, ctx, x, y, w, h, title, body, dark: bool = False):
    shape = slide.shapes.add_shape(ctx["MSO_SHAPE"].ROUNDED_RECTANGLE, ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx["RGBColor"].from_string(GRAPHITE if dark else WHITE)
    shape.line.color.rgb = ctx["RGBColor"].from_string(BRAND_GREEN if dark else "D6DED9")
    _text(slide, ctx, x + 0.16, y + 0.15, w - 0.32, 0.25, title, size=8.5, color="DCE4DF" if dark else BRAND_GREEN, bold=True)
    _text(slide, ctx, x + 0.16, y + 0.46, w - 0.32, h - 0.55, body, size=10.5, color=WHITE if dark else BRAND_DARK, bold=False)


def _step(slide, ctx, x, y, num, title, body):
    _text(slide, ctx, x, y, 0.62, 0.42, num, size=18, color=BRAND_GREEN, bold=True)
    _text(slide, ctx, x, y + 0.55, 1.9, 0.35, title, size=12, color=BRAND_DARK, bold=True)
    _text(slide, ctx, x, y + 0.95, 1.95, 1.25, body, size=9, color=MUTED)


def _bullets(slide, ctx, x, y, w, h, items):
    body = "\n".join(f"• {item}" for item in items)
    _text(slide, ctx, x, y, w, h, body, size=11, color=BRAND_DARK)


def _table(slide, ctx, x, y, w, h, rows):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ctx["RGBColor"].from_string(BRAND_DARK if r == 0 else WHITE)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = ctx["Pt"](8.5 if r else 8)
                paragraph.font.bold = r == 0
                paragraph.font.color.rgb = ctx["RGBColor"].from_string(WHITE if r == 0 else BRAND_DARK)


def _chart(slide, ctx, x, y, w, h, title, categories, series, chart_type="column"):
    chart_data = ctx["CategoryChartData"]()
    chart_data.categories = [str(item) for item in categories]
    for name, values in series:
        chart_data.add_series(name, [float(value or 0) for value in values])
    chart_map = {
        "column": ctx["XL_CHART_TYPE"].COLUMN_CLUSTERED,
        "line": ctx["XL_CHART_TYPE"].LINE_MARKERS,
        "pie": ctx["XL_CHART_TYPE"].PIE,
        "bar": ctx["XL_CHART_TYPE"].BAR_CLUSTERED,
    }
    chart_shape = slide.shapes.add_chart(
        chart_map.get(chart_type, ctx["XL_CHART_TYPE"].COLUMN_CLUSTERED),
        ctx["Inches"](x),
        ctx["Inches"](y),
        ctx["Inches"](w),
        ctx["Inches"](h),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = chart_type != "pie"
    if chart.has_legend:
        chart.legend.position = ctx["XL_LEGEND_POSITION"].BOTTOM
        chart.legend.include_in_layout = False
    if chart_type == "pie":
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.position = ctx["XL_LABEL_POSITION"].BEST_FIT
        chart.plots[0].data_labels.show_percentage = True
    return chart


def _text(slide, ctx, x, y, w, h, text, size=12, color=BRAND_DARK, bold=False):
    box = slide.shapes.add_textbox(ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.font.name = "Arial"
    paragraph.font.size = ctx["Pt"](size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = ctx["RGBColor"].from_string(color)
    return box


def _footer(slide, ctx, label):
    _text(slide, ctx, 0.65, 7.05, 7.6, 0.24, label, size=7.5, color=MUTED)
    _text(slide, ctx, 11.6, 7.05, 1.0, 0.24, "MG", size=7.5, color=BRAND_GREEN, bold=True)


def _summary_metrics(financials: dict) -> list[tuple[str, str]]:
    periods = financials.get("periods") or ["FY2025"]
    latest = periods[-1]
    revenue = _line_value(financials, "income_statement", "Revenue", latest)
    ebitda = _line_value(financials, "income_statement", "EBITDA", latest)
    debt = _line_value(financials, "balance_sheet", "Debt", latest)
    margin = (ebitda / revenue) if revenue else 0
    return [
        ("Latest period", latest),
        ("Revenue", _money(revenue)),
        ("EBITDA margin", f"{margin:.1%}"),
        ("Debt", _money(debt)),
    ]


def _historical_table(financials: dict) -> list[list[str]]:
    periods = (financials.get("periods") or ["FY2023", "FY2024", "FY2025"])[-3:]
    rows = [["Metric"] + periods]
    for section, metric in [("income_statement", "Revenue"), ("income_statement", "EBITDA"), ("balance_sheet", "Cash"), ("balance_sheet", "Debt")]:
        rows.append([metric] + [_money(_line_value(financials, section, metric, period)) for period in periods])
    return rows


def _bp_projection_series(ctx) -> tuple[list[str], list[float], list[float], list[float]]:
    assumptions = ctx.get("assumptions", {})
    revenue_streams = assumptions.get("revenue_streams") or []
    cost_base = assumptions.get("cost_base") or {}
    debt_tranches = assumptions.get("debt_tranches") or []
    years = ["Y1", "Y2", "Y3", "Y4", "Y5"]
    revenue = []
    ebitda = []
    debt = []
    opening_debt = sum(float(item.get("opening_balance") or 0) for item in debt_tranches) or 500000
    cogs = float(cost_base.get("cogs_percent") or 0.35)
    fixed_opex = float(cost_base.get("opex_fixed_monthly") or 80000) * 12
    for year in range(5):
        annual_revenue = 0
        for item in revenue_streams:
            monthly = float(item.get("volume") or 0) * float(item.get("price") or 0)
            growth = (float(item.get("volume_growth") or 0) + float(item.get("price_growth") or 0)) * 12
            annual_revenue += monthly * 12 * ((1 + growth) ** year)
        if annual_revenue == 0:
            annual_revenue = 1200000 * (1.05 ** year)
        revenue.append(annual_revenue)
        ebitda.append(annual_revenue * (1 - cogs) - fixed_opex * (1.03 ** year))
        debt.append(max(0, opening_debt * (1 - 0.12 * year)))
    return years, revenue, ebitda, debt


def _line_value(financials: dict, section: str, name: str, period: str) -> float:
    for line in financials.get(section, []):
        if str(line.get("name", "")).lower() == name.lower():
            return float((line.get("values") or {}).get(period) or 0)
    return 0.0


def _money(value: float) -> str:
    abs_value = abs(value)
    if abs_value >= 1_000_000:
        return f"{value / 1_000_000:.1f}m"
    if abs_value >= 1_000:
        return f"{value / 1_000:.1f}k"
    return f"{value:,.0f}"
